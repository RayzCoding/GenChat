#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
✔ Supports new JSON schema (type / content / url / fontLimit)
✔ Preserves template font styles
✔ Group shape support
✔ image field: replace only when URL is provided; keep template image otherwise
✔ shape.name matches JSON key
✔ Per-slide background support (full-screen background image)
✔ Automatic template slide cleanup
✔ duplicate_slide copies image relationships completely
"""

import argparse
import json
import os
import sys
from pptx import Presentation
from urllib.parse import urlparse
from io import BytesIO
import requests
from copy import deepcopy

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

# =========================
# Utility functions
# =========================

def is_url(text):
    if not isinstance(text, str):
        return False
    try:
        r = urlparse(text)
        return r.scheme and r.netloc
    except:
        return False


def download_image(url):
    try:
        r = requests.get(url, timeout=20)
        r.raise_for_status()
        return BytesIO(r.content)
    except Exception as e:
        print(f"❌ Image download failed: {url}, error: {e}")
        return None


# =========================
# Iterate all shapes (supports groups)
# =========================

def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            for sub in iter_shapes(shape.shapes):
                yield sub


# =========================
# Replace text while preserving styles (supports fontLimit)
# =========================
def replace_text_keep_style(shape, value, font_limit=None):
    if not shape.has_text_frame:
        return

    text = "" if value is None else str(value)

    # --- Truncation logic ---
    if font_limit:
        # Truncate when text exceeds the limit by 3+ characters
        # Keep full text when only 1-2 characters over the limit
        if len(text) > (font_limit + 2):
            text = text[:font_limit]
    # -----------------------

    tf = shape.text_frame

    if not tf.paragraphs:
        tf.text = text
        return

    p = tf.paragraphs[0]

    if not p.runs:
        p.add_run().text = text
    else:
        p.runs[0].text = text
        # Remove extra runs
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)

    # Remove extra paragraphs
    for extra in tf.paragraphs[1:]:
        extra._element.getparent().remove(extra._element)


# =========================
# Set slide background
# =========================

def set_slide_background(slide, url):
    """
    Set PPT background image (true slide background, not a shape).
    """

    img_data = download_image(url)
    if not img_data:
        return False

    try:
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import qn

        # 1️⃣ Add image to relationships
        image_part, rId = slide.part.get_or_add_image_part(img_data)

        # 2️⃣ Remove old background
        slide_elm = slide._element
        for bg in slide_elm.findall(qn("p:bg")):
            slide_elm.remove(bg)

        # 3️⃣ Create background XML
        bg_xml = f"""
        <p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <p:bgPr>
            <a:blipFill dpi="0" rotWithShape="1">
              <a:blip r:embed="{rId}"/>
              <a:stretch>
                <a:fillRect/>
              </a:stretch>
            </a:blipFill>
          </p:bgPr>
        </p:bg>
        """

        bg_element = parse_xml(bg_xml)

        # 4️⃣ Insert background
        slide_elm.insert(0, bg_element)

        print("✓ Background image set successfully")
        return True

    except Exception as e:
        print("❌ Failed to set background:", e)
        return False


# =========================
# ⭐ Duplicate slide
# =========================

def duplicate_slide(prs, index):
    """
    Duplicate a slide (copy background node and reuse media resources to avoid file bloat).
    """
    source_slide = prs.slides[index]
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # ========= 1. Clear auto-generated placeholder shapes on the new slide
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # ========= 2. Copy relationships and map old rId to new rId
    rId_mapping = {}
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        try:
            if rel.is_external:
                # External links use target_ref directly
                new_rel = new_slide.part.rels.get_or_add_extRel(rel.reltype, rel.target_ref)
                rId_mapping[rel.rId] = new_rel.rId
            else:
                # Internal resources (images, etc.) must bind to the original target_part
                # so the new slide references existing media without doubling file size
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_mapping[rel.rId] = new_rId
        except Exception as e:
            pass

    def update_xml_rids(element):
        """Walk XML nodes and replace old rIds with new rIds."""
        for el in element.iter():
            for attr_name, attr_value in list(el.attrib.items()):
                if attr_value in rId_mapping:
                    el.set(attr_name, rId_mapping[attr_value])

    # ========= 3. Copy shapes
    from copy import deepcopy
    for shape in source_slide.shapes:
        new_el = deepcopy(shape._element)
        update_xml_rids(new_el)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # ========= 4. Copy slide-specific background
    # Background node <p:bg> is nested under <p:cSld>!
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    source_bg = source_slide._element.cSld.find(f"{{{ns}}}bg")

    if source_bg is not None:
        # Remove old background on the new slide
        new_bg = new_slide._element.cSld.find(f"{{{ns}}}bg")
        if new_bg is not None:
            new_slide._element.cSld.remove(new_bg)

        # Copy background XML and update rIds inside
        bg_copy = deepcopy(source_bg)
        update_xml_rids(bg_copy)

        # Insert background at the front of <p:cSld>
        new_slide._element.cSld.insert(0, bg_copy)

    return new_slide


# =========================
# Fill slide (supports new JSON schema)
# =========================

def fill_slide(slide, data):
    all_shapes = list(iter_shapes(slide.shapes))

    # name → shape
    shape_map = {
        (s.name or "").strip(): s
        for s in all_shapes if s.name
    }

    for key, field in data.items():

        if not isinstance(field, dict):
            continue

        field_type = field.get("type", "").lower()
        content = field.get("content")
        url = field.get("url")

        # Support fontLimit / font-limit
        font_limit = field.get("fontLimit") or field.get("font-limit")

        # ========= background (no shape required)
        if field_type == "background":
            if is_url(url):
                set_slide_background(slide, url)
            continue

        shape = shape_map.get(key)
        if not shape:
            print(f"⚠️ Shape not found: {key}")
            continue

        # ========= image
        if field_type == "image":
            if is_url(url):
                img_data = download_image(url)
                if img_data:
                    try:
                        left, top, width, height = (
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height
                        )

                        parent = shape._element.getparent()
                        parent.remove(shape._element)

                        slide.shapes.add_picture(
                            img_data,
                            left,
                            top,
                            width=width,
                            height=height
                        )

                        print(f"✓ Image replaced: {key}")

                    except Exception as e:
                        print(f"❌ Image replacement failed {key}: {e}")

            # Keep template image when URL is not provided
            continue

        # ========= text
        if field_type == "text":
            replace_text_keep_style(
                shape,
                content,
                font_limit
            )


# =========================
# Delete slide
# =========================

def delete_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    prs.part.drop_rel(slides[index].rId)
    del slide_id_list[index]


# =========================
# Main render flow
# =========================

def render_ppt(template_path, schema_json, output_path):
    print("🚀 Starting render flow...")

    try:
        schema = json.loads(schema_json)
    except Exception as e:
        print(f"❌ JSON parse failed: {e}")
        return

    slides_data = schema.get("slides", [])
    prs = Presentation(template_path)
    template_count = len(prs.slides)

    print(f"📊 Template slide count: {template_count}")

    # Generate new slides from JSON
    for i, slide_def in enumerate(slides_data):
        t_idx = slide_def.get("templatePageIndex", 1) - 1
        data = slide_def.get("data", {})

        if 0 <= t_idx < template_count:
            print(f"📝 Rendering slide {i+1} (template {t_idx+1})")
            new_slide = duplicate_slide(prs, t_idx)
            fill_slide(new_slide, data)
        else:
            print(f"⚠️ Template index out of range: {t_idx+1}")

    # Delete template slides
    for i in reversed(range(template_count)):
        delete_slide(prs, i)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)

    print("✅ Render completed:", output_path)


# =========================
# CLI
# =========================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    schema_json = os.environ.get("PPT_SCHEMA")

    if not schema_json:
        schema_file = os.environ.get("PPT_SCHEMA_FILE")
        if schema_file and os.path.exists(schema_file):
            with open(schema_file, "r", encoding="utf-8") as f:
                schema_json = f.read()

    if not schema_json:
        print("❌ PPT_SCHEMA or PPT_SCHEMA_FILE not provided")
        sys.exit(1)

    render_ppt(args.template, schema_json, args.output)


if __name__ == "__main__":
    main()
